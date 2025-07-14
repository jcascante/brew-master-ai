#!/usr/bin/env python3
"""
Unified processing engine for Brew Master AI data processing.
Combines all processing features: audio extraction, transcription, text processing, embeddings, and cleanup.
"""

import os
import subprocess
import shutil
import re
import json
import hashlib
import time
import logging
from typing import List, Dict, Any, Optional, Tuple, Set
from dataclasses import dataclass
from pathlib import Path
from datetime import datetime
import unicodedata

# Data processing imports
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
import spacy
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.http import models as qmodels

# Audio and image processing
import whisper
from pptx import Presentation
import pytesseract
from PIL import Image

# Import our unified config
from config import Config, ConfigManager

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


@dataclass
class ProcessingResult:
    """Result of a processing operation"""
    success: bool
    files_processed: int
    files_skipped: int
    files_failed: int
    total_time: float
    output_files: List[str]
    errors: List[str]
    metadata: Dict[str, Any]


@dataclass
class CleanupResult:
    """Result of cleanup operation"""
    files_checked: int
    files_orphaned: int
    chunks_deleted: int
    files_cleaned: List[str]
    errors: List[str]


class DataValidator:
    """Validates and cleans text data"""
    
    def __init__(self, config: Config):
        self.config = config
        self._setup_nlp()
    
    def _setup_nlp(self):
        """Setup NLP components"""
        try:
            # Download required NLTK data
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
            nltk.download('wordnet', quiet=True)
            nltk.download('averaged_perceptron_tagger', quiet=True)
            
            self.stop_words = set(stopwords.words(self.config.preprocessing.language))
            self.lemmatizer = WordNetLemmatizer()
            
            # Try to load spaCy model for better sentence segmentation
            try:
                self.nlp = spacy.load("en_core_web_sm")
            except OSError:
                logger.warning("spaCy model not found. Using NLTK for sentence tokenization.")
                self.nlp = None
                
        except Exception as e:
            logger.error(f"Error setting up NLP components: {e}")
            self.stop_words = set()
            self.lemmatizer = None
            self.nlp = None
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""
        
        # Unicode normalization
        if self.config.preprocessing.normalize_unicode:
            text = unicodedata.normalize('NFKC', text)
        
        # Remove special characters
        if self.config.preprocessing.remove_special_chars:
            text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\{\}]', '', text)
        
        # Convert to lowercase
        if self.config.preprocessing.lowercase:
            text = text.lower()
        
        # Remove numbers
        if self.config.preprocessing.remove_numbers:
            text = re.sub(r'\d+', '', text)
        
        # Remove punctuation
        if self.config.preprocessing.remove_punctuation:
            text = re.sub(r'[^\w\s]', '', text)
        
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    def validate_text(self, text: str) -> Tuple[bool, str]:
        """Validate text quality"""
        if not text:
            return False, "Empty text"
        
        # Check minimum length
        if len(text) < self.config.validation.min_text_length:
            return False, f"Text too short ({len(text)} chars, min {self.config.validation.min_text_length})"
        
        # Check maximum length
        if len(text) > self.config.validation.max_text_length:
            return False, f"Text too long ({len(text)} chars, max {self.config.validation.max_text_length})"
        
        # Check for meaningful content
        words = text.split()
        if len(words) < 5:
            return False, "Not enough meaningful words"
        
        # Check for repetitive content
        unique_words = set(words)
        if len(unique_words) / len(words) < 0.15:
            return False, "Too much repetitive content"
        
        return True, "Valid text"
    
    def preprocess_text(self, text: str) -> str:
        """Preprocess text according to configuration"""
        if not self.config.preprocessing.clean_text:
            return text
        
        # Clean text
        text = self.clean_text(text)
        
        # Remove stopwords
        if self.config.preprocessing.remove_stopwords:
            words = text.split()
            words = [word for word in words if word.lower() not in self.stop_words]
            text = ' '.join(words)
        
        # Lemmatization
        if self.config.preprocessing.lemmatize and self.lemmatizer:
            words = text.split()
            words = [self.lemmatizer.lemmatize(word) for word in words]
            text = ' '.join(words)
        
        return text


class TextChunker:
    """Chunks text into smaller pieces for embedding"""
    
    def __init__(self, config: Config):
        self.config = config
        self._setup_nlp()
    
    def _setup_nlp(self):
        """Setup NLP components for sentence tokenization"""
        try:
            nltk.download('punkt', quiet=True)
            if spacy.util.is_package("en_core_web_sm"):
                self.nlp = spacy.load("en_core_web_sm")
            else:
                self.nlp = None
                logger.warning("spaCy model not available, using NLTK for sentence tokenization")
        except Exception as e:
            logger.warning(f"Could not setup NLP: {e}")
            self.nlp = None
    
    def chunk_text(self, text: str, metadata: Dict[str, Any]) -> List[Tuple[str, Dict[str, Any]]]:
        """Chunk text into smaller pieces"""
        if self.config.text_processing.chunk_by_sentences:
            return self._chunk_by_sentences(text, metadata)
        else:
            return self._chunk_by_size(text, metadata)
    
    def _chunk_by_sentences(self, text: str, metadata: Dict[str, Any]) -> List[Tuple[str, Dict[str, Any]]]:
        """Chunk text by sentences with overlap"""
        # Tokenize into sentences
        if self.nlp:
            doc = self.nlp(text)
            sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
        else:
            sentences = sent_tokenize(text)
        
        if not sentences:
            return []
        
        chunks = []
        chunk_index = 0
        
        # Create chunks with sentence boundaries
        current_chunk = []
        current_size = 0
        
        for i, sentence in enumerate(sentences):
            sentence_size = len(sentence)
            
            # Check if adding this sentence would exceed max size
            if (current_size + sentence_size > self.config.text_processing.max_chunk_size and 
                current_chunk and len(current_chunk) >= self.config.text_processing.max_sentences_per_chunk):
                
                # Create chunk
                chunk_text = ' '.join(current_chunk)
                if len(chunk_text) >= self.config.text_processing.min_chunk_size:
                    chunk_metadata = self._create_chunk_metadata(metadata, chunk_index, i - len(current_chunk), i - 1)
                    chunks.append((chunk_text, chunk_metadata))
                    chunk_index += 1
                
                # Start new chunk with overlap
                if self.config.text_processing.overlap_size > 0:
                    overlap_sentences = self._get_overlap_sentences(current_chunk, self.config.text_processing.overlap_size)
                    current_chunk = overlap_sentences
                    current_size = sum(len(s) for s in current_chunk)
                else:
                    current_chunk = []
                    current_size = 0
            
            current_chunk.append(sentence)
            current_size += sentence_size
        
        # Add final chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            if len(chunk_text) >= self.config.text_processing.min_chunk_size:
                chunk_metadata = self._create_chunk_metadata(metadata, chunk_index, len(sentences) - len(current_chunk), len(sentences) - 1)
                chunks.append((chunk_text, chunk_metadata))
        
        return chunks
    
    def _chunk_by_size(self, text: str, metadata: Dict[str, Any]) -> List[Tuple[str, Dict[str, Any]]]:
        """Chunk text by character size with overlap"""
        chunks = []
        chunk_index = 0
        
        for i in range(0, len(text), self.config.text_processing.max_chunk_size - self.config.text_processing.overlap_size):
            chunk_text = text[i:i + self.config.text_processing.max_chunk_size]
            
            if len(chunk_text) >= self.config.text_processing.min_chunk_size:
                chunk_metadata = self._create_chunk_metadata(metadata, chunk_index, i, i + len(chunk_text))
                chunks.append((chunk_text, chunk_metadata))
                chunk_index += 1
        
        return chunks
    
    def _get_overlap_sentences(self, sentences: List[str], overlap_size: int) -> List[str]:
        """Get sentences for overlap based on size"""
        overlap_sentences = []
        current_size = 0
        
        for sentence in reversed(sentences):
            if current_size + len(sentence) <= overlap_size:
                overlap_sentences.insert(0, sentence)
                current_size += len(sentence)
            else:
                break
        
        return overlap_sentences
    
    def _create_chunk_metadata(self, base_metadata: Dict[str, Any], chunk_index: int, 
                              start_pos: int, end_pos: int) -> Dict[str, Any]:
        """Create metadata for a chunk"""
        metadata = base_metadata.copy()
        metadata.update({
            'chunk_index': chunk_index,
            'start_position': start_pos,
            'end_position': end_pos,
            'chunk_size': end_pos - start_pos,
            'processing_timestamp': datetime.now().isoformat()
        })
        return metadata


class MetadataEnricher:
    """Enriches metadata with file and content information"""
    
    def __init__(self):
        self.brewing_keywords = {
            'beer', 'brew', 'brewing', 'malt', 'hops', 'yeast', 'fermentation',
            'wort', 'mash', 'boil', 'lager', 'ale', 'stout', 'ipa', 'pilsner',
            'barley', 'wheat', 'rye', 'oats', 'cascade', 'citra', 'mosaic'
        }
    
    def enrich_metadata(self, file_path: str, content_type: str, text: str) -> Dict[str, Any]:
        """Enrich metadata with comprehensive information"""
        file_info = self._get_file_info(file_path)
        content_info = self._get_content_info(text, content_type)
        
        metadata = {
            'source_file': os.path.basename(file_path),
            'source_path': file_path,
            'content_type': content_type,
            'content_hash': self._generate_content_hash(text),
            **file_info,
            **content_info
        }
        
        return metadata
    
    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get file information"""
        try:
            stat = os.stat(file_path)
            return {
                'file_size': stat.st_size,
                'file_modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                'file_extension': Path(file_path).suffix.lower()
            }
        except Exception as e:
            logger.warning(f"Could not get file info for {file_path}: {e}")
            return {}
    
    def _get_content_info(self, text: str, content_type: str) -> Dict[str, Any]:
        """Get content analysis information"""
        words = text.lower().split()
        sentences = sent_tokenize(text)
        
        # Brewing-specific analysis
        found_keywords = [word for word in words if word in self.brewing_keywords]
        
        return {
            'brewing_keywords_found': len(found_keywords),
            'brewing_keywords': list(set(found_keywords)),
            'content_density': len([w for w in words if len(w) > 3]) / max(len(words), 1),
            'avg_sentence_length': len(words) / max(len(sentences), 1),
            'total_words': len(words),
            'total_sentences': len(sentences),
            'content_length': len(text)
        }
    
    def _generate_content_hash(self, text: str) -> str:
        """Generate hash for content deduplication"""
        return hashlib.md5(text.encode()).hexdigest()


class BrewMasterProcessor:
    """Main processing engine with all features"""
    
    def __init__(self, config: Config):
        self.config = config
        self.validator = DataValidator(config)
        self.chunker = TextChunker(config)
        self.metadata_enricher = MetadataEnricher()
        self.qdrant_client = QdrantClient(host=config.vector_db_host, port=config.vector_db_port)
        
        # Statistics
        self.stats = {
            'files_processed': 0,
            'chunks_created': 0,
            'chunks_validated': 0,
            'chunks_rejected': 0,
            'total_text_length': 0
        }
    
    def extract_audio(self, input_dir: str, output_dir: str) -> ProcessingResult:
        """Extract audio from video files using ffmpeg"""
        logger.info(f"Extracting audio from {input_dir} to {output_dir}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        start_time = time.time()
        files_processed = 0
        files_skipped = 0
        files_failed = 0
        output_files = []
        errors = []
        
        for filename in os.listdir(input_dir):
            if filename.lower().endswith('.mp4'):
                input_path = os.path.join(input_dir, filename)
                base = os.path.splitext(filename)[0]
                output_path = os.path.join(output_dir, base + '.wav')
                
                # Check if output already exists
                if os.path.exists(output_path):
                    logger.info(f"Skipping {filename} (audio already exists)")
                    files_skipped += 1
                    continue
                
                logger.info(f'Extracting audio from {filename}...')
                
                try:
                    subprocess.run([
                        'ffmpeg', '-y', '-i', input_path, '-vn', '-acodec', 'pcm_s16le', 
                        '-ar', str(self.config.input_processing.audio_sample_rate), 
                        '-ac', str(self.config.input_processing.audio_channels), output_path
                    ], check=True, capture_output=True)
                    
                    logger.info(f'Audio saved to {output_path}')
                    output_files.append(output_path)
                    files_processed += 1
                    
                except subprocess.CalledProcessError as e:
                    error_msg = f'Error extracting audio from {filename}: {e}'
                    logger.error(error_msg)
                    errors.append(error_msg)
                    files_failed += 1
                except Exception as e:
                    error_msg = f'Unexpected error processing {filename}: {e}'
                    logger.error(error_msg)
                    errors.append(error_msg)
                    files_failed += 1
        
        total_time = time.time() - start_time
        
        return ProcessingResult(
            success=files_failed == 0,
            files_processed=files_processed,
            files_skipped=files_skipped,
            files_failed=files_failed,
            total_time=total_time,
            output_files=output_files,
            errors=errors,
            metadata={'operation': 'audio_extraction'}
        )
    
    def transcribe_audio(self, input_dir: str, output_dir: str) -> ProcessingResult:
        """Transcribe audio files using Whisper with enhanced quality settings"""
        logger.info(f"Transcribing audio from {input_dir} to {output_dir}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        # Get all WAV files
        wav_files = [f for f in os.listdir(input_dir) if f.lower().endswith('.wav')]
        total_files = len(wav_files)
        
        if total_files == 0:
            logger.warning("No WAV files found in audio directory!")
            return ProcessingResult(
                success=False,
                files_processed=0,
                files_skipped=0,
                files_failed=0,
                total_time=0,
                output_files=[],
                errors=["No WAV files found"],
                metadata={'operation': 'transcription'}
            )
        
        logger.info(f"Found {total_files} WAV files to transcribe")
        logger.info(f"Loading Whisper model: {self.config.input_processing.whisper_model}")
        
        try:
            model = whisper.load_model(self.config.input_processing.whisper_model)
            logger.info("Whisper model loaded successfully!")
        except Exception as e:
            error_msg = f"Failed to load Whisper model: {e}"
            logger.error(error_msg)
            return ProcessingResult(
                success=False,
                files_processed=0,
                files_skipped=0,
                files_failed=total_files,
                total_time=0,
                output_files=[],
                errors=[error_msg],
                metadata={'operation': 'transcription'}
            )
        
        start_time = time.time()
        files_processed = 0
        files_skipped = 0
        files_failed = 0
        output_files = []
        errors = []
        
        # Process files one by one with progress
        for i, filename in enumerate(wav_files, 1):
            audio_path = os.path.join(input_dir, filename)
            base = os.path.splitext(filename)[0]
            transcript_path = os.path.join(output_dir, base + '.txt')
            
            # Check if transcript already exists
            if os.path.exists(transcript_path):
                logger.info(f"[{i}/{total_files}] Skipping {filename} (transcript already exists)")
                files_skipped += 1
                continue
            
            # Get file size for progress info
            file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
            
            logger.info(f"\n[{i}/{total_files}] Transcribing {filename} ({file_size_mb:.1f} MB)...")
            logger.info(f"  Input: {audio_path}")
            logger.info(f"  Output: {transcript_path}")
            
            file_start_time = time.time()
            
            try:
                logger.info("  Starting transcription with enhanced settings...")
                logger.info(f"  Using language: {self.config.input_processing.whisper_language}")
                
                # Enhanced Whisper transcription with all quality parameters
                result = model.transcribe(
                    audio_path,
                    language=self.config.input_processing.whisper_language,
                    task=self.config.input_processing.whisper_task,
                    verbose=self.config.input_processing.whisper_verbose,
                    fp16=self.config.input_processing.whisper_fp16,
                    temperature=self.config.input_processing.whisper_temperature,
                    compression_ratio_threshold=self.config.input_processing.whisper_compression_ratio_threshold,
                    logprob_threshold=self.config.input_processing.whisper_logprob_threshold,
                    no_speech_threshold=self.config.input_processing.whisper_no_speech_threshold,
                    condition_on_previous_text=self.config.input_processing.whisper_condition_on_previous_text,
                    initial_prompt=self.config.input_processing.whisper_initial_prompt,
                    best_of=self.config.input_processing.whisper_best_of,
                    beam_size=self.config.input_processing.whisper_beam_size,
                    patience=self.config.input_processing.whisper_patience,
                    length_penalty=self.config.input_processing.whisper_length_penalty,
                    suppress_tokens=self.config.input_processing.whisper_suppress_tokens,
                    suppress_blank=self.config.input_processing.whisper_suppress_blank,
                    word_timestamps=self.config.input_processing.whisper_word_timestamps,
                    prepend_punctuations=self.config.input_processing.whisper_prepend_punctuations,
                    append_punctuations=self.config.input_processing.whisper_append_punctuations
                )
                
                # Enhanced post-processing for Spanish text
                processed_text = self._post_process_spanish_text_enhanced(result['text'], result.get('segments', []))
                
                # Save transcript
                with open(transcript_path, 'w', encoding='utf-8') as f:
                    f.write(processed_text)
                
                # Calculate timing and quality metrics
                elapsed_time = time.time() - file_start_time
                words = len(result['text'].split())
                
                # Calculate average confidence if available
                avg_confidence = 0.0
                if 'segments' in result and result['segments']:
                    confidences = [seg.get('avg_logprob', 0) for seg in result['segments']]
                    avg_confidence = sum(confidences) / len(confidences) if confidences else 0.0
                
                logger.info(f"  ✅ Success! Transcript saved to {transcript_path}")
                logger.info(f"  📊 Stats: {words} words, {elapsed_time:.1f} seconds ({elapsed_time/60:.1f} minutes)")
                logger.info(f"  ⚡ Speed: {words/elapsed_time:.1f} words/second")
                if avg_confidence > 0:
                    logger.info(f"  🎯 Average confidence: {avg_confidence:.3f}")
                
                output_files.append(transcript_path)
                files_processed += 1
                
            except Exception as e:
                error_msg = f"  ❌ Error transcribing {filename}: {e}"
                logger.error(error_msg)
                errors.append(error_msg)
                files_failed += 1
        
        total_time = time.time() - start_time
        
        logger.info(f"\n🎉 Transcription complete! Processed {files_processed} files.")
        logger.info(f"Transcripts saved to: {output_dir}")
        
        return ProcessingResult(
            success=files_failed == 0,
            files_processed=files_processed,
            files_skipped=files_skipped,
            files_failed=files_failed,
            total_time=total_time,
            output_files=output_files,
            errors=errors,
            metadata={'operation': 'transcription', 'whisper_model': self.config.input_processing.whisper_model}
        )
    
    def extract_images(self, input_dir: str, output_dir: str) -> ProcessingResult:
        """Extract images from PowerPoint presentations"""
        logger.info(f"Extracting images from {input_dir} to {output_dir}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        start_time = time.time()
        files_processed = 0
        files_skipped = 0
        files_failed = 0
        output_files = []
        errors = []
        
        for filename in os.listdir(input_dir):
            if filename.lower().endswith('.pptx'):
                pptx_path = os.path.join(input_dir, filename)
                base = os.path.splitext(filename)[0]
                
                logger.info(f'Extracting images from {filename}...')
                
                try:
                    pres = Presentation(pptx_path)
                    img_count = 0
                    
                    for i, slide in enumerate(pres.slides):
                        for shape in slide.shapes:
                            if hasattr(shape, 'image'):
                                img = shape.image
                                ext = img.ext
                                img_bytes = img.blob
                                img_filename = f"{base}_slide{i+1}_img{img_count+1}.{ext}"
                                img_path = os.path.join(output_dir, img_filename)
                                
                                with open(img_path, 'wb') as f:
                                    f.write(img_bytes)
                                logger.info(f"Extracted image: {img_path}")
                                output_files.append(img_path)
                                img_count += 1
                    
                    files_processed += 1
                    
                except Exception as e:
                    error_msg = f'Error processing {filename}: {e}'
                    logger.error(error_msg)
                    errors.append(error_msg)
                    files_failed += 1
        
        total_time = time.time() - start_time
        
        return ProcessingResult(
            success=files_failed == 0,
            files_processed=files_processed,
            files_skipped=files_skipped,
            files_failed=files_failed,
            total_time=total_time,
            output_files=output_files,
            errors=errors,
            metadata={'operation': 'image_extraction'}
        )
    
    def ocr_images(self, input_dir: str, output_dir: str) -> ProcessingResult:
        """Extract text from images using OCR"""
        logger.info(f"Running OCR on {input_dir} to {output_dir}")
        
        os.makedirs(output_dir, exist_ok=True)
        
        start_time = time.time()
        files_processed = 0
        files_skipped = 0
        files_failed = 0
        output_files = []
        errors = []
        
        for filename in os.listdir(input_dir):
            if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
                image_path = os.path.join(input_dir, filename)
                base = os.path.splitext(filename)[0]
                text_path = os.path.join(output_dir, base + '.txt')
                
                # Check if output already exists
                if os.path.exists(text_path):
                    logger.info(f"Skipping {filename} (OCR text already exists)")
                    files_skipped += 1
                    continue
                
                logger.info(f'Running OCR on {filename}...')
                
                try:
                    text = pytesseract.image_to_string(Image.open(image_path))
                    with open(text_path, 'w', encoding='utf-8') as f:
                        f.write(text)
                    
                    logger.info(f'OCR text saved to {text_path}')
                    output_files.append(text_path)
                    files_processed += 1
                    
                except Exception as e:
                    error_msg = f'Error processing {filename}: {e}'
                    logger.error(error_msg)
                    errors.append(error_msg)
                    files_failed += 1
        
        total_time = time.time() - start_time
        
        return ProcessingResult(
            success=files_failed == 0,
            files_processed=files_processed,
            files_skipped=files_skipped,
            files_failed=files_failed,
            total_time=total_time,
            output_files=output_files,
            errors=errors,
            metadata={'operation': 'ocr'}
        )
    
    def process_text(self, input_dir: str, content_type: str) -> ProcessingResult:
        """Process text files with validation, chunking, and metadata enrichment"""
        logger.info(f"Processing text files in {input_dir} (content type: {content_type})")
        
        start_time = time.time()
        files_processed = 0
        files_skipped = 0
        files_failed = 0
        output_files = []
        errors = []
        all_chunks = []
        
        for filename in os.listdir(input_dir):
            if filename.lower().endswith('.txt'):
                file_path = os.path.join(input_dir, filename)
                
                logger.info(f"Processing file: {file_path}")
                
                try:
                    # Read file
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    
                    # Validate original text
                    is_valid, reason = self.validator.validate_text(text)
                    if not is_valid:
                        logger.warning(f"File {file_path} validation failed: {reason}")
                        files_failed += 1
                        continue
                    
                    # Preprocess text
                    processed_text = self.validator.preprocess_text(text)
                    
                    # Enrich metadata
                    metadata = self.metadata_enricher.enrich_metadata(file_path, content_type, processed_text)
                    
                    # Chunk text
                    chunks = self.chunker.chunk_text(processed_text, metadata)
                    
                    # Validate chunks
                    valid_chunks = []
                    for chunk_text, chunk_metadata in chunks:
                        is_valid, reason = self.validator.validate_text(chunk_text)
                        if is_valid:
                            valid_chunks.append((chunk_text, chunk_metadata))
                            self.stats['chunks_validated'] += 1
                        else:
                            logger.debug(f"Chunk rejected: {reason}")
                            self.stats['chunks_rejected'] += 1
                    
                    # Update statistics
                    self.stats['files_processed'] += 1
                    self.stats['chunks_created'] += len(chunks)
                    self.stats['total_text_length'] += len(processed_text)
                    
                    all_chunks.extend(valid_chunks)
                    output_files.append(file_path)
                    files_processed += 1
                    
                    logger.info(f"Created {len(valid_chunks)} valid chunks from {file_path}")
                    
                except Exception as e:
                    error_msg = f"Error processing file {file_path}: {e}"
                    logger.error(error_msg)
                    errors.append(error_msg)
                    files_failed += 1
        
        total_time = time.time() - start_time
        
        return ProcessingResult(
            success=files_failed == 0,
            files_processed=files_processed,
            files_skipped=files_skipped,
            files_failed=files_failed,
            total_time=total_time,
            output_files=output_files,
            errors=errors,
            metadata={
                'operation': 'text_processing',
                'content_type': content_type,
                'chunks_created': len(all_chunks),
                'chunks_validated': self.stats['chunks_validated'],
                'chunks_rejected': self.stats['chunks_rejected']
            }
        )
    
    def create_embeddings(self, input_dir: str, config_name: str = None) -> ProcessingResult:
        """Create embeddings from text files and upload to vector database"""
        logger.info(f"Creating embeddings from {input_dir}")
        
        # Get smart config if not specified
        if not config_name and self.config.smart_config:
            config_name = self.get_smart_config_for_content_type('manual')
        
        if not config_name:
            config_name = self.config.default_config
        
        logger.info(f"Using config: {config_name}")
        
        # Process text files
        text_result = self.process_text(input_dir, 'manual')
        
        if not text_result.success:
            return text_result
        
        # Generate embeddings
        logger.info("Generating embeddings...")
        
        try:
            model = SentenceTransformer(self.config.text_processing.embedding_model)
            
            # Get all chunks from the text processing
            chunks = []
            for file_path in text_result.output_files:
                # This is a simplified version - in practice, you'd need to track chunks
                # For now, we'll process the files again to get chunks
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                
                processed_text = self.validator.preprocess_text(text)
                metadata = self.metadata_enricher.enrich_metadata(file_path, 'manual', processed_text)
                file_chunks = self.chunker.chunk_text(processed_text, metadata)
                chunks.extend(file_chunks)
            
            if not chunks:
                logger.warning("No chunks to embed!")
                return ProcessingResult(
                    success=False,
                    files_processed=0,
                    files_skipped=0,
                    files_failed=0,
                    total_time=0,
                    output_files=[],
                    errors=["No chunks to embed"],
                    metadata={'operation': 'embedding_creation'}
                )
            
            # Generate embeddings
            texts = [chunk[0] for chunk in chunks]
            embeddings = model.encode(texts, show_progress_bar=True, convert_to_numpy=True)
            
            # Upload to Qdrant
            collection_name = self.config.text_processing.collection_name
            
            # Create collection if not exists
            try:
                collections = self.qdrant_client.get_collections()
                collection_exists = any(c.name == collection_name for c in collections.collections)
                
                if not collection_exists:
                    self.qdrant_client.recreate_collection(
                        collection_name=collection_name,
                        vectors_config=qmodels.VectorParams(
                            size=embeddings.shape[1], 
                            distance=self.config.text_processing.distance_metric
                        ),
                        optimizers_config=qmodels.OptimizersConfigDiff(
                            indexing_threshold=self.config.vector_db_indexing_threshold,
                            memmap_threshold=self.config.vector_db_memmap_threshold
                        )
                    )
                    logger.info(f"Created collection: {collection_name} with indexing threshold: 1000")
            except Exception as e:
                logger.error(f"Error with collection: {e}")
                return ProcessingResult(
                    success=False,
                    files_processed=0,
                    files_skipped=0,
                    files_failed=0,
                    total_time=0,
                    output_files=[],
                    errors=[f"Collection error: {e}"],
                    metadata={'operation': 'embedding_creation'}
                )
            
            # Upload points with enhanced metadata
            points = []
            for idx, (text, metadata) in enumerate(chunks):
                # Add config tracking
                enhanced_metadata = metadata.copy()
                enhanced_metadata['config_used'] = config_name
                enhanced_metadata['file_path'] = metadata.get('source_path', '')
                enhanced_metadata['processing_timestamp'] = datetime.now().isoformat()
                
                points.append(qmodels.PointStruct(
                    id=idx,
                    vector=embeddings[idx].tolist(),
                    payload=enhanced_metadata | {"text": text}
                ))
            
            self.qdrant_client.upsert(collection_name=collection_name, wait=True, points=points)
            
            logger.info(f"Successfully uploaded {len(points)} embeddings to Qdrant")
            
            return ProcessingResult(
                success=True,
                files_processed=text_result.files_processed,
                files_skipped=text_result.files_skipped,
                files_failed=text_result.files_failed,
                total_time=text_result.total_time,
                output_files=text_result.output_files,
                errors=text_result.errors,
                metadata={
                    'operation': 'embedding_creation',
                    'config_used': config_name,
                    'embeddings_uploaded': len(points),
                    'collection_name': collection_name
                }
            )
            
        except Exception as e:
            error_msg = f"Error creating embeddings: {e}"
            logger.error(error_msg)
            return ProcessingResult(
                success=False,
                files_processed=0,
                files_skipped=0,
                files_failed=0,
                total_time=0,
                output_files=[],
                errors=[error_msg],
                metadata={'operation': 'embedding_creation'}
            )
    
    def get_smart_config_for_content_type(self, content_type: str) -> str:
        """Automatically select appropriate config based on content type"""
        return self.config.content_type_configs.get(content_type, self.config.default_config)
    
    def get_existing_file_chunks(self) -> Dict[str, List[int]]:
        """Get existing chunks grouped by source file"""
        try:
            points = self.qdrant_client.scroll(
                collection_name=self.config.text_processing.collection_name,
                limit=10000,
                with_payload=True
            )[0]
            
            file_chunks = {}
            for point in points:
                source_file = point.payload.get('source_file', '')
                if source_file:
                    if source_file not in file_chunks:
                        file_chunks[source_file] = []
                    file_chunks[source_file].append(point.id)
            
            logger.info(f"Found {len(file_chunks)} files with existing chunks")
            return file_chunks
            
        except Exception as e:
            logger.warning(f"Could not retrieve existing chunks: {e}")
            return {}
    
    def get_processed_files_configs(self) -> Dict[str, str]:
        """Get which config was used to process each file"""
        try:
            points = self.qdrant_client.scroll(
                collection_name=self.config.text_processing.collection_name,
                limit=10000,
                with_payload=True
            )[0]
            
            file_configs = {}
            for point in points:
                source_file = point.payload.get('source_file', '')
                config_used = point.payload.get('config_used', 'unknown')
                if source_file and config_used:
                    file_configs[source_file] = config_used
            
            logger.info(f"Found {len(file_configs)} files with config tracking")
            return file_configs
            
        except Exception as e:
            logger.warning(f"Could not retrieve file configs: {e}")
            return {}
    
    def cleanup_orphaned_chunks(self, data_dirs: List[str]) -> CleanupResult:
        """Remove chunks from files that no longer exist"""
        logger.info("Starting cleanup of orphaned chunks...")
        
        # Get existing chunks grouped by file
        existing_file_chunks = self.get_existing_file_chunks()
        if not existing_file_chunks:
            logger.info("No existing chunks found, skipping cleanup")
            return CleanupResult(
                files_checked=0,
                files_orphaned=0,
                chunks_deleted=0,
                files_cleaned=[],
                errors=[]
            )
        
        # Get current files
        current_files = set()
        for data_dir in data_dirs:
            if os.path.exists(data_dir):
                for filename in os.listdir(data_dir):
                    if filename.lower().endswith('.txt'):
                        file_path = os.path.join(data_dir, filename)
                        current_files.add(file_path)
        
        logger.info(f"Found {len(current_files)} current files")
        
        # Find orphaned files (exist in DB but not in filesystem)
        orphaned_files = set(existing_file_chunks.keys()) - current_files
        
        cleanup_stats = CleanupResult(
            files_checked=len(existing_file_chunks),
            files_orphaned=len(orphaned_files),
            chunks_deleted=0,
            files_cleaned=[],
            errors=[]
        )
        
        if not orphaned_files:
            logger.info("No orphaned files found")
            return cleanup_stats
        
        # Delete chunks from orphaned files
        for orphaned_file in orphaned_files:
            chunk_ids = existing_file_chunks[orphaned_file]
            try:
                self.qdrant_client.delete(
                    collection_name=self.config.text_processing.collection_name,
                    points_selector=qmodels.PointIdsList(points=chunk_ids)
                )
                cleanup_stats.chunks_deleted += len(chunk_ids)
                cleanup_stats.files_cleaned.append(orphaned_file)
                logger.info(f"Deleted {len(chunk_ids)} chunks from orphaned file: {orphaned_file}")
                
            except Exception as e:
                error_msg = f"Error deleting chunks from {orphaned_file}: {e}"
                logger.error(error_msg)
                cleanup_stats.errors.append(error_msg)
        
        logger.info(f"Cleanup completed: {cleanup_stats.chunks_deleted} chunks deleted from {len(orphaned_files)} files")
        return cleanup_stats
    
    def should_process_file(self, file_path: str, config_name: str) -> bool:
        """Check if file should be processed (avoid duplicates)"""
        if not self.config.cleanup.deduplication:
            return True
        
        processed_configs = self.get_processed_files_configs()
        filename = os.path.basename(file_path)
        
        if filename in processed_configs:
            existing_config = processed_configs[filename]
            if existing_config == config_name:
                logger.info(f"Skipping {filename} - already processed with {config_name}")
                return False
            else:
                logger.info(f"Reprocessing {filename} - config changed from {existing_config} to {config_name}")
                return True
        
        logger.info(f"Processing {filename} - new file")
        return True
    
    def get_statistics(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return self.stats.copy()
    
    def _post_process_spanish_text(self, text: str) -> str:
        """Post-process Spanish text to improve transcription quality"""
        if not text:
            return text
        
        # Common Spanish transcription fixes
        replacements = {
            # Common abbreviations and contractions
            r'\bq\b': 'que',
            r'\bx\b': 'por',
            r'\bd\b': 'de',
            r'\bpa\b': 'para',
            r'\bta\b': 'también',
            r'\bke\b': 'que',
            
            # Brewing terminology fixes
            r'\blupulo\b': 'lúpulo',
            r'\bmalta\b': 'malta',
            r'\blevadura\b': 'levadura',
            r'\bfermentacion\b': 'fermentación',
            r'\bmacerado\b': 'macerado',
            r'\bhervor\b': 'hervor',
            r'\bcarbonatacion\b': 'carbonatación',
            r'\bembotellado\b': 'embotellado',
            r'\bmaduracion\b': 'maduración',
            
            # Common Spanish word fixes
            r'\besta\b': 'está',
            r'\beste\b': 'este',
            r'\bpara\b': 'para',
            r'\bpor\b': 'por',
            r'\bcon\b': 'con',
            r'\bdel\b': 'del',
            r'\bal\b': 'al',
            
            # Punctuation fixes
            r'\s+([.,!?;:])': r'\1',  # Remove spaces before punctuation
            r'([.,!?;:])\s*([A-Z])': r'\1 \2',  # Add space after punctuation before capital letters
        }
        
        # Apply replacements
        for pattern, replacement in replacements.items():
            text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
        
        # Fix capitalization
        sentences = text.split('. ')
        fixed_sentences = []
        for sentence in sentences:
            if sentence:
                # Capitalize first letter of sentence
                sentence = sentence[0].upper() + sentence[1:] if sentence else sentence
                fixed_sentences.append(sentence)
        
        text = '. '.join(fixed_sentences)
        
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text


# Backward compatibility functions
def create_enhanced_embeddings(config: Config = None):
    """Create embeddings using the enhanced processing pipeline (backward compatibility)"""
    if config is None:
        config = Config()
    
    processor = BrewMasterProcessor(config)
    
    # Process different content types
    all_chunks = []
    
    # Process transcripts
    if os.path.exists('data/transcripts'):
        logger.info("Processing video transcripts...")
        result = processor.process_text('data/transcripts', 'transcript')
        if result.success:
            all_chunks.extend(result.output_files)
    
    # Process OCR text
    if os.path.exists('data/presentation_texts'):
        logger.info("Processing presentation OCR text...")
        result = processor.process_text('data/presentation_texts', 'ocr')
        if result.success:
            all_chunks.extend(result.output_files)
    
    if not all_chunks:
        logger.warning("No text files found to process!")
        return
    
    # Create embeddings
    result = processor.create_embeddings('data/transcripts', 'general_brewing')
    
    # Log statistics
    stats = processor.get_statistics()
    logger.info("Processing Statistics:")
    for key, value in stats.items():
        logger.info(f"  {key}: {value}")


if __name__ == "__main__":
    # Test the processor
    print("Testing Brew Master Processor...")
    
    config = Config()
    processor = BrewMasterProcessor(config)
    
    # Test configuration
    print(f"Default config: {config.default_config}")
    print(f"Smart config: {config.smart_config}")
    print(f"Whisper model: {config.input_processing.whisper_model}")
    print(f"Chunk size: {config.text_processing.max_chunk_size}")
    
    # Test smart config selection
    smart_config = processor.get_smart_config_for_content_type('transcript')
    print(f"Smart config for transcript: {smart_config}")
    
    print("Processor test completed!") 