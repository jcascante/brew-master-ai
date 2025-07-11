#!/usr/bin/env python3
"""
Enhanced Brew Master AI Data Extraction CLI
Features advanced chunking strategies, data validation, and quality assessment.
"""

import argparse
import os
import sys
from pathlib import Path

# Import our enhanced modules
from enhanced_processor import create_enhanced_embeddings, ProcessingConfig
from chunking_configs import get_config, list_available_configs, create_custom_config
from data_validator import DataQualityAnalyzer
from config_loader import ConfigLoader, create_argument_parser

def extract_audio(input_dir, output_dir, processed_dir=None):
    """Extract audio from video files using ffmpeg"""
    import subprocess
    
    os.makedirs(output_dir, exist_ok=True)
    
    for filename in os.listdir(input_dir):
        if filename.lower().endswith('.mp4'):
            input_path = os.path.join(input_dir, filename)
            base = os.path.splitext(filename)[0]
            output_path = os.path.join(output_dir, base + '.wav')
            print(f'Extracting audio from {filename}...')
            
            try:
                subprocess.run([
                    'ffmpeg', '-y', '-i', input_path, '-vn', '-acodec', 'pcm_s16le', 
                    '-ar', '16000', '-ac', '1', output_path
                ], check=True, capture_output=True)
                print(f'Audio saved to {output_path}')
                
            except subprocess.CalledProcessError as e:
                print(f'Error extracting audio from {filename}: {e}')
            except Exception as e:
                print(f'Unexpected error processing {filename}: {e}')

def transcribe_audio(audio_dir, transcript_dir, processed_audio_dir=None):
    """Transcribe audio files using Whisper"""
    import whisper
    import time
    from pathlib import Path
    
    os.makedirs(transcript_dir, exist_ok=True)
    
    # Get all WAV files
    wav_files = [f for f in os.listdir(audio_dir) if f.lower().endswith('.wav')]
    total_files = len(wav_files)
    
    if total_files == 0:
        print("No WAV files found in audio directory!")
        return
    
    print(f"Found {total_files} WAV files to transcribe")
    print("Loading Whisper model...")
    model = whisper.load_model('base')
    print("Model loaded successfully!")
    
    # Process files one by one with progress
    for i, filename in enumerate(wav_files, 1):
        audio_path = os.path.join(audio_dir, filename)
        base = os.path.splitext(filename)[0]
        transcript_path = os.path.join(transcript_dir, base + '.txt')
        
        # Check if transcript already exists
        if os.path.exists(transcript_path):
            print(f"[{i}/{total_files}] Skipping {filename} (transcript already exists)")
            continue
        
        # Get file size for progress info
        file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
        
        print(f"\n[{i}/{total_files}] Transcribing {filename} ({file_size_mb:.1f} MB)...")
        print(f"  Input: {audio_path}")
        print(f"  Output: {transcript_path}")
        
        start_time = time.time()
        
        try:
            print("  Starting transcription...")
            result = model.transcribe(audio_path)
            
            # Save transcript
            with open(transcript_path, 'w', encoding='utf-8') as f:
                f.write(result['text'])
            
            # Calculate timing
            elapsed_time = time.time() - start_time
            words = len(result['text'].split())
            
            print(f"  ✅ Success! Transcript saved to {transcript_path}")
            print(f"  📊 Stats: {words} words, {elapsed_time:.1f} seconds ({elapsed_time/60:.1f} minutes)")
            print(f"  ⚡ Speed: {words/elapsed_time:.1f} words/second")
            
        except Exception as e:
            print(f"  ❌ Error transcribing {filename}: {e}")
            # Continue with next file instead of stopping
    
    print(f"\n🎉 Transcription complete! Processed {total_files} files.")
    print(f"Transcripts saved to: {transcript_dir}")

def extract_pptx_images(pptx_dir, images_dir, processed_pptx_dir=None):
    """Extract images from PowerPoint presentations"""
    from pptx import Presentation
    
    os.makedirs(images_dir, exist_ok=True)
    
    for filename in os.listdir(pptx_dir):
        if filename.lower().endswith('.pptx'):
            pptx_path = os.path.join(pptx_dir, filename)
            base = os.path.splitext(filename)[0]
            
            print(f'Extracting images from {filename}...')
            try:
                pres = Presentation(pptx_path)
                img_count = 0
                
                for i, slide in enumerate(pres.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, 'image'):
                            img = shape.image
                            ext = img.ext
                            img_bytes = img.blob
                            img_filename = f"{base}_slide{i+1}_img{img_count+1}.{ext}"
                            img_path = os.path.join(images_dir, img_filename)
                            
                            with open(img_path, 'wb') as f:
                                f.write(img_bytes)
                            print(f"Extracted image: {img_path}")
                            img_count += 1
                
            except Exception as e:
                print(f'Error processing {filename}: {e}')

def ocr_images(images_dir, texts_dir, processed_images_dir=None):
    """Extract text from images using OCR"""
    import pytesseract
    from PIL import Image
    
    os.makedirs(texts_dir, exist_ok=True)
    
    for filename in os.listdir(images_dir):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            image_path = os.path.join(images_dir, filename)
            base = os.path.splitext(filename)[0]
            text_path = os.path.join(texts_dir, base + '.txt')
            
            print(f'Running OCR on {filename}...')
            try:
                text = pytesseract.image_to_string(Image.open(image_path))
                with open(text_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                print(f'OCR text saved to {text_path}')
                
            except Exception as e:
                print(f'Error processing {filename}: {e}')

def validate_data(directory_path, output_report=None, create_plots=False):
    """Validate and analyze data quality"""
    print(f"Validating data in: {directory_path}")
    
    analyzer = DataQualityAnalyzer()
    results = analyzer.analyze_directory(directory_path)
    
    # Generate report
    report = analyzer.generate_report(results, output_report)
    print(report)
    
    # Create visualizations if requested
    if create_plots:
        try:
            analyzer.create_visualizations(results, "quality_plots")
            print("Visualizations created in 'quality_plots' directory")
        except Exception as e:
            print(f"Could not create visualizations: {e}")
    
    return results

def main():
    # Create argument parser with config override options
    parser = create_argument_parser()
    
    # Add our specific arguments
    parser.add_argument('--extract-audio', action='store_true', 
                       help='Extract audio from videos')
    parser.add_argument('--transcribe-audio', action='store_true', 
                       help='Transcribe audio to text')
    parser.add_argument('--extract-pptx-images', action='store_true', 
                       help='Extract images from PowerPoint presentations')
    parser.add_argument('--ocr-images', action='store_true', 
                       help='Extract text from images using OCR')
    
    # Enhanced embedding creation
    parser.add_argument('--create-embeddings', action='store_true', 
                       help='Create enhanced embeddings and upload to Qdrant')
    
    # Data validation
    parser.add_argument('--validate', type=str, metavar='DIRECTORY',
                       help='Validate data quality in specified directory')
    parser.add_argument('--report', type=str, metavar='FILE',
                       help='Output file for quality report')
    parser.add_argument('--plots', action='store_true',
                       help='Create quality visualization plots')
    
    # Configuration management
    parser.add_argument('--list-configs', action='store_true',
                       help='List available chunking configurations')
    

    
    args = parser.parse_args()
    
    # Load configuration (YAML file + command-line overrides)
    loader = ConfigLoader()
    config = loader.load_config(args)
    
    # List configurations
    if args.list_configs:
        list_available_configs()
        return
    
    # Validate data
    if args.validate:
        validate_data(args.validate, args.report, args.plots)
        return
    
    # Data extraction pipeline
    if args.extract_audio:
        print(f"Extracting audio from {config.videos_dir} to {config.audios_dir}")
        extract_audio(config.videos_dir, config.audios_dir)
    
    if args.transcribe_audio:
        print(f"Transcribing audio from {config.audios_dir} to {config.transcripts_dir}")
        transcribe_audio(config.audios_dir, config.transcripts_dir)
        if args.validate:
            print("\nValidating transcripts...")
            validate_data(config.transcripts_dir, args.report, args.plots)
    
    if args.extract_pptx_images:
        print(f"Extracting images from {config.presentations_dir} to {config.images_dir}")
        extract_pptx_images(config.presentations_dir, config.images_dir)
    
    if args.ocr_images:
        print(f"Running OCR on {config.images_dir} to {config.presentation_texts_dir}")
        ocr_images(config.images_dir, config.presentation_texts_dir)
        if args.validate:
            print("\nValidating OCR text...")
            validate_data(config.presentation_texts_dir, args.report, args.plots)
    
    # Create enhanced embeddings
    if args.create_embeddings:
        print(f"Creating embeddings with configuration: {args.config}")
        
        # Get configuration
        if args.chunk_size or args.overlap or args.min_chunk or args.max_sentences:
            # Use custom configuration
            config = create_custom_config(
                max_chunk_size=args.chunk_size or 1000,
                overlap_size=args.overlap or 200,
                min_chunk_size=args.min_chunk or 150,
                max_sentences_per_chunk=args.max_sentences or 10
            )
            print("Using custom chunking configuration")
        else:
            # Use preset configuration
            config = get_config(args.config)
            print(f"Using preset configuration: {args.config}")
        
        create_enhanced_embeddings(config)
    
    # If no arguments provided, show help
    if not any(vars(args).values()):
        parser.print_help()

if __name__ == '__main__':
    main() 