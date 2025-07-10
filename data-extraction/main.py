import argparse
import os
import subprocess
import shutil


def extract_audio(input_dir, output_dir, processed_dir):
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(processed_dir, exist_ok=True)
    for filename in os.listdir(input_dir):
        if filename.lower().endswith('.mp4'):
            input_path = os.path.join(input_dir, filename)
            base = os.path.splitext(filename)[0]
            output_path = os.path.join(output_dir, base + '.wav')
            print(f'Extracting audio from {filename}...')
            subprocess.run([
                'ffmpeg', '-y', '-i', input_path, '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1', output_path
            ], check=True)
            print(f'Audio saved to {output_path}')
            # Move processed video
            processed_path = os.path.join(processed_dir, filename)
            shutil.move(input_path, processed_path)
            print(f'Moved {filename} to {processed_dir}')


def transcribe_audio(audio_dir, transcript_dir, processed_audio_dir):
    import whisper
    os.makedirs(transcript_dir, exist_ok=True)
    os.makedirs(processed_audio_dir, exist_ok=True)
    model = whisper.load_model('base')
    for filename in os.listdir(audio_dir):
        if filename.lower().endswith('.wav'):
            audio_path = os.path.join(audio_dir, filename)
            base = os.path.splitext(filename)[0]
            transcript_path = os.path.join(transcript_dir, base + '.txt')
            print(f'Transcribing {filename}...')
            result = model.transcribe(audio_path)
            with open(transcript_path, 'w') as f:
                f.write(result['text'])
            print(f'Transcript saved to {transcript_path}')
            # Move processed audio
            processed_path = os.path.join(processed_audio_dir, filename)
            shutil.move(audio_path, processed_path)
            print(f'Moved {filename} to {processed_audio_dir}')


def extract_pptx_images(pptx_dir, images_dir, processed_pptx_dir):
    from pptx import Presentation
    os.makedirs(images_dir, exist_ok=True)
    os.makedirs(processed_pptx_dir, exist_ok=True)
    for filename in os.listdir(pptx_dir):
        if filename.lower().endswith('.pptx'):
            pptx_path = os.path.join(pptx_dir, filename)
            base = os.path.splitext(filename)[0]
            pres = Presentation(pptx_path)
            img_count = 0
            for i, slide in enumerate(pres.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        img = shape.image
                        ext = img.ext
                        img_bytes = img.blob
                        img_filename = f"{base}_slide{i+1}_img{img_count+1}.{ext}"
                        img_path = os.path.join(images_dir, img_filename)
                        with open(img_path, 'wb') as f:
                            f.write(img_bytes)
                        print(f"Extracted image: {img_path}")
                        img_count += 1
            # Move processed pptx
            processed_path = os.path.join(processed_pptx_dir, filename)
            shutil.move(pptx_path, processed_path)
            print(f"Moved {filename} to {processed_pptx_dir}")


def ocr_images(images_dir, texts_dir, processed_images_dir):
    import pytesseract
    from PIL import Image
    os.makedirs(texts_dir, exist_ok=True)
    os.makedirs(processed_images_dir, exist_ok=True)
    for filename in os.listdir(images_dir):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif')):
            image_path = os.path.join(images_dir, filename)
            base = os.path.splitext(filename)[0]
            text_path = os.path.join(texts_dir, base + '.txt')
            print(f'Running OCR on {filename}...')
            text = pytesseract.image_to_string(Image.open(image_path))
            with open(text_path, 'w') as f:
                f.write(text)
            print(f'OCR text saved to {text_path}')
            # Move processed image
            processed_path = os.path.join(processed_images_dir, filename)
            shutil.move(image_path, processed_path)
            print(f'Moved {filename} to {processed_images_dir}')


def create_embeddings_and_upload():
    from sentence_transformers import SentenceTransformer
    from qdrant_client import QdrantClient
    from qdrant_client.http import models as qmodels
    import glob

    # Gather all text files
    transcript_files = glob.glob('data/transcripts/*.txt')
    ocr_files = glob.glob('data/presentation_texts/*.txt')
    all_files = transcript_files + ocr_files

    # Read and chunk text
    docs = []
    metadatas = []
    chunk_size = 500
    for file_path in all_files:
        with open(file_path, 'r') as f:
            text = f.read()
        # Chunk text
        for i in range(0, len(text), chunk_size):
            chunk = text[i:i+chunk_size]
            docs.append(chunk)
            metadatas.append({
                'source_file': os.path.basename(file_path),
                'chunk_index': i // chunk_size
            })

    print(f"Loaded {len(docs)} text chunks from {len(all_files)} files.")

    # Generate embeddings
    model = SentenceTransformer('all-MiniLM-L6-v2')
    embeddings = model.encode(docs, show_progress_bar=True, convert_to_numpy=True)

    # Connect to Qdrant
    client = QdrantClient(host="localhost", port=6333)
    collection_name = "brew_master_ai"

    # Create collection if not exists
    if collection_name not in [c.name for c in client.get_collections().collections]:
        client.recreate_collection(
            collection_name=collection_name,
            vectors_config=qmodels.VectorParams(size=embeddings.shape[1], distance="Cosine")
        )

    # Upload points
    points = []
    for idx, (embedding, metadata) in enumerate(zip(embeddings, metadatas)):
        points.append(qmodels.PointStruct(
            id=idx,
            vector=embedding.tolist(),
            payload=metadata | {"text": docs[idx]}
        ))
    client.upsert(collection_name=collection_name, points=points)
    print(f"Uploaded {len(points)} embeddings to Qdrant collection '{collection_name}'.")


def main():
    parser = argparse.ArgumentParser(description='Brew Master AI Data Extraction CLI')
    parser.add_argument('--hello', action='store_true', help='Print hello world')
    parser.add_argument('--extract-audio', action='store_true', help='Extract audio from videos')
    parser.add_argument('--transcribe-audio', action='store_true', help='Transcribe audio to text')
    parser.add_argument('--extract-pptx-images', action='store_true', help='Extract images from PowerPoint presentations')
    parser.add_argument('--ocr-images', action='store_true', help='Extract text from images using OCR')
    parser.add_argument('--create-embeddings', action='store_true', help='Create embeddings and upload to Qdrant')
    args = parser.parse_args()

    if args.hello:
        print('Hello from Data Extraction CLI!')
    if args.extract_audio:
        extract_audio('data/videos', 'data/audios', 'data/processed')
    if args.transcribe_audio:
        transcribe_audio('data/audios', 'data/transcripts', 'data/processed_audios')
    if args.extract_pptx_images:
        extract_pptx_images('data/presentations', 'data/presentation_images', 'data/processed_presentations')
    if args.ocr_images:
        ocr_images('data/presentation_images', 'data/presentation_texts', 'data/processed_images')
    if args.create_embeddings:
        create_embeddings_and_upload()

if __name__ == '__main__':
    main()
